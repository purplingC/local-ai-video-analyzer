from concurrent import futures
import grpc
import textwrap
from pathlib import Path
from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.enum.text import PP_ALIGN
from reportlab.lib.pagesizes import letter
from reportlab.pdfgen import canvas
from transformers import pipeline
from grpc_services import video_analysis_pb2, video_analysis_pb2_grpc

ARTIFACTS_DIR = Path(__file__).resolve().parents[1] / "artifacts"
UPLOADS_DIR = Path(__file__).resolve().parents[1] / "uploads"
ARTIFACTS_DIR.mkdir(exist_ok=True)

# Initialize summarization model
try:
    summarizer = pipeline("summarization", model="models/t5-small", tokenizer="models/t5-small")
    _HAS_SUMMARY = True
except Exception:
    summarizer = None
    _HAS_SUMMARY = False


def _stub(port: int):
    ch = grpc.insecure_channel(f"localhost:{port}")
    return ch, video_analysis_pb2_grpc.VideoAnalysisStub(ch)


class GenerationServicer(video_analysis_pb2_grpc.VideoAnalysisServicer):
    def GenerateReport(self, request, context):
        file_path = request.file_path
        report_type = (request.report_type or "pdf").lower()
        base = Path(file_path).stem

        transcript_path = UPLOADS_DIR / f"{base}.txt"
        vision_path = UPLOADS_DIR / f"{base}.vision.txt"

        # Auto-call agents if required files are missing
        if not transcript_path.exists():
            print("Transcript not found — auto-calling Transcription Agent...")
            ch, stub = _stub(50051)
            stub.TranscribeVideo(video_analysis_pb2.VideoRequest(file_path=file_path))
            ch.close()

        if not vision_path.exists():
            print("Vision results not found — auto-calling Vision Agent...")
            ch, stub = _stub(50052)
            stub.AnalyzeVideo(video_analysis_pb2.VideoRequest(file_path=file_path))
            ch.close()

        transcript_text = transcript_path.read_text(encoding="utf-8") if transcript_path.exists() else ""
        vision_summary = vision_path.read_text(encoding="utf-8") if vision_path.exists() else ""

        # Summarize transcript if possible
        if _HAS_SUMMARY and transcript_text:
            short_summary = summarizer(
                transcript_text[:1000],
                max_length=150,
                min_length=40,
                do_sample=False
            )[0]["summary_text"]
        else:
            short_summary = transcript_text[:800] or "No transcript available."

        # Format vision summary with bullet points
        if "Objects detected:" in vision_summary:
            parts = vision_summary.split("Objects detected:")
            objects_list = parts[1].strip().splitlines() if len(parts) > 1 else []
            bullet_list = "\n".join([f"• {obj.strip()}" for obj in objects_list if obj.strip()])
            formatted_vision = f"Objects detected:\n{bullet_list}"
        else:
            formatted_vision = vision_summary.strip() or "No visual data available."

        # Section labels
        title_text = "Local AI Video Analyzer — Summary Report"
        transcript_header = "Transcript Summary"
        vision_header = "Vision Summary"

        # Generate PowerPoint Report
        if report_type in ("pptx", "ppt"):
            out_path = ARTIFACTS_DIR / f"{base}_summary.pptx"
            prs = Presentation()
            slide_layout = prs.slide_layouts[6]  # blank slide
            slide = prs.slides.add_slide(slide_layout)

            # Add centered title
            title_box = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(8), Inches(1))
            title_tf = title_box.text_frame
            title_p = title_tf.add_paragraph()
            title_p.text = title_text
            title_p.font.size = Pt(26)
            title_p.font.bold = True
            title_p.font.name = "Arial"
            title_p.alignment = PP_ALIGN.CENTER

            # Add content box (Transcript + Vision)
            content_box = slide.shapes.add_textbox(Inches(1), Inches(1.8), Inches(8.5), Inches(5))
            content_tf = content_box.text_frame

            # Transcript section
            p1 = content_tf.add_paragraph()
            p1.text = transcript_header
            p1.font.bold = True
            p1.font.size = Pt(16)
            p1.font.name = "Arial"
            p1.space_after = Pt(6)

            p2 = content_tf.add_paragraph()
            p2.text = short_summary.strip()
            p2.font.size = Pt(14)
            p2.font.name = "Arial"
            p2.space_after = Pt(10)

            # Vision section
            p3 = content_tf.add_paragraph()
            p3.text = vision_header
            p3.font.bold = True
            p3.font.size = Pt(16)
            p3.font.name = "Arial"
            p3.space_after = Pt(6)

            # Vision bullet list
            for line in formatted_vision.splitlines():
                if line.startswith("Objects detected"):
                    continue  # skip repeating header line
                if line.startswith("•"):
                    p = content_tf.add_paragraph()
                    p.text = line
                    p.font.size = Pt(14)
                    p.font.name = "Arial"
                    p.space_after = Pt(4)

            prs.save(str(out_path))
            print(f"PowerPoint summary saved: {out_path}")
            return video_analysis_pb2.ReportResponse(report_path=str(out_path))

        # Generate PDF Report 
        out_path = ARTIFACTS_DIR / f"{base}_summary.pdf"
        c = canvas.Canvas(str(out_path), pagesize=letter)
        width, height = letter

        # Title centered
        c.setFont("Helvetica-Bold", 20)
        title_width = c.stringWidth(title_text, "Helvetica-Bold", 20)
        c.drawString((width - title_width) / 2, 760, title_text)

        # Separator line
        c.setLineWidth(1)
        c.line(40, 755, width - 40, 755)

        y = 730

        # Transcript Section
        c.setFont("Helvetica-Bold", 14)
        c.drawString(40, y, transcript_header)
        y -= 20
        c.setFont("Helvetica", 12)
        for line in textwrap.wrap(short_summary, width=90):
            c.drawString(40, y, line)
            y -= 15

        # Vision Section
        y -= 25
        c.setFont("Helvetica-Bold", 14)
        c.drawString(40, y, vision_header)
        y -= 20
        c.setFont("Helvetica", 12)

        for line in formatted_vision.splitlines():
            if line.startswith("Objects detected"):
                c.drawString(40, y, "Objects detected:")
                y -= 15
            elif line.startswith("•"):
                c.drawString(60, y, line)
                y -= 15

        c.showPage()
        c.save()
        print(f"PDF summary saved: {out_path}")
        return video_analysis_pb2.ReportResponse(report_path=str(out_path))


def serve():
    server = grpc.server(futures.ThreadPoolExecutor(max_workers=2))
    video_analysis_pb2_grpc.add_VideoAnalysisServicer_to_server(GenerationServicer(), server)
    server.add_insecure_port("[::]:50053")
    msg = "Generation Agent running on port 50053"
    msg += " (Summarization enabled)" if _HAS_SUMMARY else " (no summarizer)"
    print(msg)
    server.start()
    server.wait_for_termination()


if __name__ == "__main__":
    serve()
